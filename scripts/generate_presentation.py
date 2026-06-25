#!/usr/bin/env python3
"""Generate a professional PowerPoint for the Multi-Agent Research System project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "docs" / "Multi_Agent_Research_System.pptx"

# Professional palette — navy + Google-inspired accents
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
NAVY_MID = RGBColor(0x1A, 0x3A, 0x5C)
ACCENT = RGBColor(0x42, 0x85, 0xF4)
ACCENT_LIGHT = RGBColor(0xE8, 0xF0, 0xFE)
TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
BODY = RGBColor(0x42, 0x42, 0x42)
MUTED = RGBColor(0x75, 0x75, 0x75)
ROW_ALT = RGBColor(0xF5, 0xF7, 0xFA)
BORDER = RGBColor(0xDE, 0xE2, 0xE8)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.65)
MARGIN_R = Inches(0.65)
CONTENT_TOP = Inches(1.55)
FOOTER_Y = Inches(7.05)

_slide_counter = 0
PROJECT_NAME = "Multi-Agent Research Intelligence System"


class Theme:
  @staticmethod
  def rect(slide, left, top, width, height, fill, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if not line:
      s.line.fill.background()
    return s

  @staticmethod
  def set_text(tf, text, size=14, bold=False, color=BODY, align=PP_ALIGN.LEFT, spacing=6):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(spacing)
    return p

  @staticmethod
  def add_footer(slide, section: str = "") -> None:
    global _slide_counter
    _slide_counter += 1
    Theme.rect(slide, Inches(0), FOOTER_Y, SLIDE_W, Inches(0.45), ACCENT_LIGHT)
    left = slide.shapes.add_textbox(MARGIN_L, FOOTER_Y + Inches(0.06), Inches(6), Inches(0.3))
    Theme.set_text(left.text_frame, PROJECT_NAME, size=9, color=MUTED)
    if section:
      mid = slide.shapes.add_textbox(Inches(4), FOOTER_Y + Inches(0.06), Inches(5.3), Inches(0.3))
      Theme.set_text(mid.text_frame, section, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    right = slide.shapes.add_textbox(Inches(11.5), FOOTER_Y + Inches(0.06), Inches(1.2), Inches(0.3))
    Theme.set_text(right.text_frame, str(_slide_counter), size=9, color=MUTED, align=PP_ALIGN.RIGHT)

  @staticmethod
  def slide_header(slide, title: str, subtitle: str = "") -> None:
    Theme.rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    Theme.rect(slide, Inches(0), Inches(0.08), Inches(0.12), Inches(1.35), NAVY)
    box = slide.shapes.add_textbox(MARGIN_L, Inches(0.35), Inches(11.5), Inches(0.65))
    Theme.set_text(box.text_frame, title, size=26, bold=True, color=NAVY)
    if subtitle:
      sub = slide.shapes.add_textbox(MARGIN_L, Inches(0.95), Inches(11.5), Inches(0.4))
      Theme.set_text(sub.text_frame, subtitle, size=12, color=MUTED)


def new_slide(prs: Presentation):
  return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
  Theme.rect(slide, Inches(0), Inches(4.85), SLIDE_W, Inches(0.06), ACCENT)
  Theme.rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)

  badge = slide.shapes.add_textbox(MARGIN_L, Inches(1.2), Inches(5), Inches(0.35))
  p = badge.text_frame.paragraphs[0]
  p.text = "CAPSTONE PROJECT  ·  GOOGLE CLOUD PLATFORM"
  p.font.size = Pt(11)
  p.font.bold = True
  p.font.color.rgb = ACCENT

  title = slide.shapes.add_textbox(MARGIN_L, Inches(1.75), Inches(11), Inches(1.8))
  tf = title.text_frame
  tf.word_wrap = True
  p = tf.paragraphs[0]
  p.text = "Multi-Agent Research\nIntelligence System"
  p.font.size = Pt(40)
  p.font.bold = True
  p.font.color.rgb = WHITE
  p.line_spacing = 1.1

  sub = slide.shapes.add_textbox(MARGIN_L, Inches(3.7), Inches(10), Inches(1.2))
  tf2 = sub.text_frame
  lines = [
    "Google ADK  ·  RAG  ·  Tavily Web Search  ·  Cloud Run",
    "",
    "Dinaz Sanaullah",
    "June 2025",
  ]
  for i, line in enumerate(lines):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(16 if i < 2 else 14)
    p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD) if i < 2 else RGBColor(0x99, 0xAA, 0xBB)
    p.space_after = Pt(4)

  link = slide.shapes.add_textbox(MARGIN_L, Inches(6.55), Inches(10), Inches(0.4))
  Theme.set_text(link.text_frame, "github.com/Dinaz-Sanaullah/multi-agent-research-system", size=11, color=ACCENT)


def add_section_slide(prs: Presentation, number: str, title: str, subtitle: str) -> None:
  slide = new_slide(prs)
  Theme.rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY_MID)
  Theme.rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)

  num = slide.shapes.add_textbox(MARGIN_L, Inches(2.2), Inches(3), Inches(1.2))
  Theme.set_text(num.text_frame, number, size=72, bold=True, color=ACCENT)

  t = slide.shapes.add_textbox(MARGIN_L, Inches(3.5), Inches(10), Inches(1))
  Theme.set_text(t.text_frame, title, size=34, bold=True, color=WHITE)

  s = slide.shapes.add_textbox(MARGIN_L, Inches(4.55), Inches(10), Inches(0.6))
  Theme.set_text(s.text_frame, subtitle, size=16, color=RGBColor(0xBB, 0xCC, 0xDD))


def add_agenda_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, "Agenda", "Overview of today's presentation")
  items = [
    ("01", "Problem & Objectives"),
    ("02", "System Architecture"),
    ("03", "Multi-Agent Design"),
    ("04", "RAG & Web Search"),
    ("05", "Deployment & Results"),
    ("06", "Challenges & Conclusion"),
  ]
  x_positions = [MARGIN_L, Inches(4.5), Inches(8.35)]
  for i, (num, label) in enumerate(items):
    col, row = i % 3, i // 3
    left = x_positions[col]
    top = CONTENT_TOP + Inches(row * 2.1)
    card = Theme.rect(slide, left, top, Inches(3.55), Inches(1.75), WHITE, line=True)
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    inner = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), Inches(3), Inches(1.3))
    tf = inner.text_frame
    p1 = tf.paragraphs[0]
    p1.text = num
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(6)
  Theme.add_footer(slide, "Agenda")


def add_bullet_slide(prs: Presentation, title: str, subtitle: str, bullets: list[str], section: str) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, title, subtitle)
  body = slide.shapes.add_textbox(MARGIN_L + Inches(0.15), CONTENT_TOP, Inches(11.8), Inches(5.2))
  tf = body.text_frame
  tf.word_wrap = True
  for i, bullet in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    clean = bullet.lstrip("• ").lstrip("✅ ").strip()
    p.text = f"▸  {clean}"
    p.font.size = Pt(15)
    p.font.color.rgb = BODY
    p.space_after = Pt(10)
    p.level = 1 if bullet.startswith("  ") else 0
  Theme.add_footer(slide, section)


def add_two_column_slide(
  prs: Presentation, title: str, subtitle: str,
  left_title: str, left_items: list[str],
  right_title: str, right_items: list[str], section: str,
) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, title, subtitle)

  def column(left, col_title, items):
    Theme.rect(slide, left, CONTENT_TOP, Inches(5.85), Inches(5.15), ACCENT_LIGHT)
    hdr = slide.shapes.add_textbox(left + Inches(0.2), CONTENT_TOP + Inches(0.15), Inches(5.4), Inches(0.4))
    Theme.set_text(hdr.text_frame, col_title, size=14, bold=True, color=NAVY)
    body = slide.shapes.add_textbox(left + Inches(0.2), CONTENT_TOP + Inches(0.55), Inches(5.4), Inches(4.4))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
      p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
      p.text = f"▸  {item}"
      p.font.size = Pt(13)
      p.font.color.rgb = BODY
      p.space_after = Pt(7)

  column(MARGIN_L, left_title, left_items)
  column(Inches(6.85), right_title, right_items)
  Theme.add_footer(slide, section)


def add_architecture_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, "System Architecture", "End-to-end request flow through the agent pipeline")

  steps = [
    ("User Query", NAVY),
    ("Orchestrator", NAVY_MID),
    ("Planner", ACCENT),
    ("RAG Researcher", TEAL),
    ("Web Researcher", TEAL),
    ("Synthesizer", ACCENT),
    ("Reviewer / QA", NAVY_MID),
    ("Response", NAVY),
  ]
  top = CONTENT_TOP + Inches(0.3)
  box_w = Inches(1.35)
  gap = Inches(0.18)
  total_w = len(steps) * box_w + (len(steps) - 1) * gap
  start_x = (SLIDE_W - total_w) / 2

  for i, (label, color) in enumerate(steps):
    left = start_x + i * (box_w + gap)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, Inches(0.85))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if i < len(steps) - 1:
      arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + box_w + Inches(0.02), top + Inches(0.3), gap - Inches(0.04), Inches(0.25))
      arrow.fill.solid()
      arrow.fill.fore_color.rgb = BORDER
      arrow.line.fill.background()

  # Pattern labels
  patterns = [
    ("Hierarchical Delegation", "Orchestrator → Pipeline", MARGIN_L),
    ("Sequential Flow", "Planner → Research → QA Loop", Inches(4.6)),
    ("Parallel / Sequential", "RAG + Web Researchers", Inches(8.75)),
  ]
  for label, desc, left in patterns:
    card = Theme.rect(slide, left, CONTENT_TOP + Inches(1.5), Inches(3.9), Inches(1.35), WHITE, line=True)
    card.line.color.rgb = BORDER
    t = slide.shapes.add_textbox(left + Inches(0.15), CONTENT_TOP + Inches(1.62), Inches(3.6), Inches(1.1))
    tf = t.text_frame
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(4)

  # Tech stack bar
  stack = slide.shapes.add_textbox(MARGIN_L, CONTENT_TOP + Inches(3.2), Inches(12), Inches(0.8))
  tf = stack.text_frame
  p = tf.paragraphs[0]
  p.text = "Technology Stack"
  p.font.size = Pt(12)
  p.font.bold = True
  p.font.color.rgb = NAVY
  p2 = tf.add_paragraph()
  p2.text = "Gemini 2.5 Flash  ·  text-embedding-005  ·  FAISS  ·  Tavily  ·  Google ADK  ·  Cloud Run"
  p2.font.size = Pt(13)
  p2.font.color.rgb = ACCENT
  p2.space_before = Pt(6)

  Theme.add_footer(slide, "Architecture")


def add_table_slide(prs: Presentation, title: str, subtitle: str, headers: list[str], rows: list[list[str]], section: str) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, title, subtitle)
  cols = len(headers)
  row_count = len(rows) + 1
  table_shape = slide.shapes.add_table(row_count, cols, MARGIN_L, CONTENT_TOP, Inches(12), Inches(0.42 * row_count))
  table = table_shape.table

  col_widths = [Inches(2.2), Inches(4.5), Inches(5.3)] if cols == 3 else None
  if col_widths:
    for i, w in enumerate(col_widths):
      table.columns[i].width = w

  for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
      p.font.bold = True
      p.font.size = Pt(11)
      p.font.color.rgb = WHITE

  for i, row in enumerate(rows, 1):
    for j, val in enumerate(row):
      cell = table.cell(i, j)
      cell.text = val
      if i % 2 == 0:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_ALT
      for p in cell.text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = BODY

  Theme.add_footer(slide, section)


def add_metrics_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, "Test Results & Performance", "Validated through RAG smoke tests, web search, and end-to-end pipeline runs")
  metrics = [
    ("104", "Chunks Indexed", "3 arXiv papers in FAISS"),
    ("0.74", "Avg RAG Score", "Similarity on test queries"),
    ("~6s", "Web Search", "Tavily 3-result query"),
    ("~57s", "Full Pipeline", "Multi-agent end-to-end"),
  ]
  for i, (value, label, detail) in enumerate(metrics):
    left = MARGIN_L + Inches(i * 3.15)
    top = CONTENT_TOP + Inches(0.4)
    card = Theme.rect(slide, left, top, Inches(2.85), Inches(2.2), WHITE, line=True)
    card.line.color.rgb = BORDER
    v = slide.shapes.add_textbox(left, top + Inches(0.35), Inches(2.85), Inches(0.9))
    Theme.set_text(v.text_frame, value, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    l = slide.shapes.add_textbox(left, top + Inches(1.2), Inches(2.85), Inches(0.4))
    Theme.set_text(l.text_frame, label, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    d = slide.shapes.add_textbox(left, top + Inches(1.6), Inches(2.85), Inches(0.45))
    Theme.set_text(d.text_frame, detail, size=10, color=MUTED, align=PP_ALIGN.CENTER)

  findings = slide.shapes.add_textbox(MARGIN_L, CONTENT_TOP + Inches(3.1), Inches(12), Inches(2))
  tf = findings.text_frame
  bullets = [
    "RAG correctly retrieved Attention, BERT, and GPT-3 paper excerpts",
    "Web search returned current transformer research via Tavily API",
    "Quality review loop produced cited, academically-toned synthesis",
    "Deployed live at research-agent-chfrwn7sia-uc.a.run.app",
  ]
  for i, b in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"▸  {b}"
    p.font.size = Pt(14)
    p.font.color.rgb = BODY
    p.space_after = Pt(8)
  Theme.add_footer(slide, "Results")


def add_challenges_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.slide_header(slide, "Challenges & Mitigations", "Production issues encountered and resolved during development")
  challenges = [
    ("429 RESOURCE_EXHAUSTED", "Gemini quota limits on multi-agent burst calls", "Retry config, sequential research mode, reduced loop iterations"),
    ("Model Not Found", "gemini-2.0-flash unavailable in project region", "Migrated to gemini-2.5-flash on Vertex AI"),
    ("Tool Registration Error", "LLM called wrong tool on researcher agents", "Both search tools registered on both agents"),
    ("Embedding Token Limit", "Batch exceeded 20K token embedding cap", "Batched embeddings (8 chunks per API call)"),
    ("Cloud Run 403", "Service required authentication for public access", "IAM invoker binding + ADK UI redeploy"),
  ]
  headers = ["Challenge", "Description", "Solution"]
  table_shape = slide.shapes.add_table(len(challenges) + 1, 3, MARGIN_L, CONTENT_TOP, Inches(12), Inches(3.2))
  table = table_shape.table
  table.columns[0].width = Inches(2.8)
  table.columns[1].width = Inches(4.5)
  table.columns[2].width = Inches(4.7)
  for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(10)
    cell.text_frame.paragraphs[0].font.color.rgb = WHITE
  for i, row in enumerate(challenges, 1):
    for j, val in enumerate(row):
      cell = table.cell(i, j)
      cell.text = val
      if i % 2 == 0:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROW_ALT
      cell.text_frame.paragraphs[0].font.size = Pt(9)
      cell.text_frame.paragraphs[0].font.color.rgb = BODY
  Theme.add_footer(slide, "Challenges")


def add_closing_slide(prs: Presentation) -> None:
  slide = new_slide(prs)
  Theme.rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
  Theme.rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)

  t = slide.shapes.add_textbox(MARGIN_L, Inches(2.4), Inches(10), Inches(1))
  Theme.set_text(t.text_frame, "Thank You", size=44, bold=True, color=WHITE)

  q = slide.shapes.add_textbox(MARGIN_L, Inches(3.5), Inches(10), Inches(0.5))
  Theme.set_text(q.text_frame, "Questions & Discussion", size=18, color=RGBColor(0xBB, 0xCC, 0xDD))

  info = slide.shapes.add_textbox(MARGIN_L, Inches(4.5), Inches(11), Inches(1.8))
  tf = info.text_frame
  lines = [
    ("Live Demo", "https://research-agent-chfrwn7sia-uc.a.run.app/"),
    ("GitHub", "github.com/Dinaz-Sanaullah/multi-agent-research-system"),
    ("GCP Project", "multi-agent-research-system"),
  ]
  for i, (label, value) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{label}:  {value}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE if i == 0 else RGBColor(0xAA, 0xCC, 0xFF)
    p.space_after = Pt(8)


def main() -> None:
  global _slide_counter
  _slide_counter = 0

  prs = Presentation()
  prs.slide_width = SLIDE_W
  prs.slide_height = SLIDE_H

  add_title_slide(prs)
  add_agenda_slide(prs)

  add_section_slide(prs, "01", "Problem & Objectives", "Defining the research assistant use case")
  add_bullet_slide(prs, "Problem Statement", "Why a multi-agent approach is needed", [
    "Researchers require answers combining established literature with the latest publications",
    "Single-prompt LLMs hallucinate citations and cannot query private paper corpora",
    "Academic knowledge bases and real-time web sources need different retrieval strategies",
    "Automated quality review is essential before presenting research findings",
  ], "Problem")

  add_bullet_slide(prs, "Project Objectives", "All seven core requirements delivered", [
    "Multi-agent system with Google ADK — Orchestrator, Researcher, Reviewer/QA",
    "RAG pipeline — PDF extraction, semantic chunking, text-embedding-005, FAISS",
    "Web search via Tavily API with intelligent RAG vs. Web vs. Hybrid routing",
    "Custom tools demonstrating agent composition and state management",
    "Observability — structured JSON logging, metrics, optional Cloud Trace",
    "Streamlit UI for local interaction and PDF ingestion",
    "Production deployment on Google Cloud Run with CI/CD",
  ], "Objectives")

  add_section_slide(prs, "02", "System Architecture", "How agents, tools, and services connect")
  add_architecture_slide(prs)

  add_section_slide(prs, "03", "Multi-Agent Design", "Seven specialized agents, four communication patterns")
  add_table_slide(prs, "Agent Roster", "Specialized roles within the research pipeline",
    ["Agent", "Responsibility", "Primary Tools"],
    [
      ["Orchestrator", "Root coordinator; delegates pipeline", "get_system_metrics"],
      ["Planner", "Query classification & research plan", "save_research_plan"],
      ["RAG Researcher", "Academic knowledge base search", "search_knowledge_base"],
      ["Web Researcher", "Real-time citation discovery", "search_web"],
      ["Synthesizer", "Merge findings into response", "format_citations"],
      ["Reviewer / QA", "Quality gate & approval", "exit_loop, append_to_state"],
      ["Refiner", "Revise synthesis on feedback", "format_citations"],
    ], "Agents")

  add_two_column_slide(prs, "Communication Patterns", "Four ADK workflow patterns implemented",
    "Workflow Patterns", [
      "Sequential — Planner → Research → QA Loop",
      "Parallel — RAG + Web researchers concurrently",
      "Hierarchical — Orchestrator delegates to pipeline",
      "Feedback Loop — Review → Refine until approved",
    ],
    "State Management", [
      "Shared session state across all agents",
      "user_query, research_plan, rag_findings",
      "web_findings, research_synthesis",
      "review_feedback via output_key templating",
    ], "Patterns")

  add_section_slide(prs, "04", "RAG & Web Search", "Dual-source intelligence with smart routing")
  add_two_column_slide(prs, "Knowledge Retrieval", "RAG pipeline and web search integration",
    "RAG Pipeline", [
      "PDF extraction with PyPDF",
      "Semantic paragraph-aware chunking",
      "Vertex AI text-embedding-005",
      "FAISS vector store (104 chunks)",
      "Top-k similarity retrieval",
    ],
    "Web Search (Tavily)", [
      "Advanced search depth, 5 results",
      "AI-generated answer summaries",
      "Authoritative source prioritization",
      "Intelligent query routing:",
      "  RAG · Web · Hybrid classification",
    ], "Retrieval")

  add_bullet_slide(prs, "Custom Tools & Observability", "Production-ready instrumentation", [
    "search_knowledge_base — FAISS retrieval over academic PDFs",
    "search_web — Tavily API for real-time information",
    "save_research_plan, append_to_state — cross-agent state sharing",
    "format_citations — merges paper and web reference sources",
    "get_system_metrics — query counts, latency, review iterations",
    "structlog JSON logging with optional OpenTelemetry Cloud Trace",
  ], "Tools")

  add_section_slide(prs, "05", "Deployment & Results", "Cloud infrastructure and validation")
  add_bullet_slide(prs, "Cloud Deployment", "Google Cloud Platform production stack", [
    "Platform: Cloud Run (us-central1) — 2 GiB RAM, 300s timeout",
    "CI/CD: Cloud Build → Docker → GCR → automated deploy",
    "Secrets: Tavily API key via Secret Manager",
    "Model: gemini-2.5-flash through Vertex AI",
    "UI: ADK Dev UI with agent graph visualization",
    "Interfaces: Streamlit (local) + ADK Web UI (production)",
  ], "Deployment")

  add_metrics_slide(prs)
  add_challenges_slide(prs)

  add_section_slide(prs, "06", "Conclusion", "Summary and future directions")
  add_bullet_slide(prs, "Key Takeaways", "Project outcomes and next steps", [
    "Delivered a production-ready multi-agent research assistant on Google ADK",
    "Combined RAG and web search with intelligent routing and QA review",
    "Deployed to Cloud Run with observability, secrets management, and CI/CD",
    "Future: persistent sessions (Firestore), larger corpora, agent evaluation benchmarks",
  ], "Conclusion")

  add_closing_slide(prs)

  OUTPUT.parent.mkdir(parents=True, exist_ok=True)
  prs.save(str(OUTPUT))
  print(f"Presentation saved to: {OUTPUT}")


if __name__ == "__main__":
  main()
